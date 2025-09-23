import mimetypes
import os
import io
import numpy as np
from PIL import Image
import cv2
import pandas as pd
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from odf.opendocument import load
from odf import text, teletype
import json
import librosa
import librosa.display
import matplotlib.pyplot as plt
from sentence_transformers import SentenceTransformer
import easyocr
import ocrmypdf
import whisper
from moviepy import VideoFileClip
import tempfile

# -------------------------------
# Initialize models
# -------------------------------
clip_model = SentenceTransformer("clip-ViT-B-32")
ocr_reader = easyocr.Reader(['en'], gpu=False)
whisper_model = whisper.load_model("small")


# -------------------------------
# Helper functions
# -------------------------------

def audio_to_spectrogram_image(file_path: str) -> Image.Image:
    y, sr = librosa.load(file_path, sr=16000, mono=True)
    S = librosa.feature.melspectrogram(y=y, sr=sr)
    S_dB = librosa.power_to_db(S, ref=np.max)
    fig, ax = plt.subplots(figsize=(3, 3), dpi=80)
    librosa.display.specshow(S_dB, sr=sr, x_axis=None, y_axis=None, ax=ax, cmap="magma")
    ax.axis("off")
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    buf.seek(0)
    return Image.open(buf).convert("RGB")


def ocr_text_from_image_easyocr(img: Image.Image, min_chars=5) -> str:
    try:
        img_np = np.array(img)
        results = ocr_reader.readtext(img_np)
        text = " ".join([res[1] for res in results])
        if len(text.strip()) >= min_chars:
            return text
    except Exception as e:
        print(f"EasyOCR error: {e}")
    return ""


def extract_text(file_path: str, mime_type: str) -> str:
    try:
        # PDF
        if mime_type == "application/pdf":
            doc = fitz.open(file_path)
            total_text_len = sum(len(page.get_text().strip()) for page in doc)
            if total_text_len < 50:
                ocr_file_path = "/tmp/ocr_output.pdf"
                ocrmypdf.ocr(file_path, ocr_file_path, force_ocr=True, language='eng')
                file_path = ocr_file_path
                doc = fitz.open(file_path)
            return "\n".join([page.get_text() for page in doc])

        # Word
        elif mime_type in ["application/msword",
                           "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            doc = docx.Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])

        # Excel / ODS
        elif mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                           "application/vnd.ms-excel",
                           "application/x-vnd.oasis.opendocument.spreadsheet"]:
            text_data = ""
            if file_path.endswith(".xlsx"):
                df = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
            elif file_path.endswith(".xls"):
                df = pd.read_excel(file_path, sheet_name=None, engine="xlrd")
            elif file_path.endswith(".ods"):
                doc = load(file_path)
                sheets = doc.spreadsheet.getElementsByType(text.Table)
                for sheet in sheets:
                    for row in sheet.getElementsByType(text.TableRow):
                        cells = row.getElementsByType(text.TableCell)
                        text_data += " | ".join([teletype.extractText(c) for c in cells]) + "\n"
                return text_data
            for sheet_name, sheet in df.items():
                text_data += f"Sheet: {sheet_name}\n"
                text_data += "\n".join(sheet.astype(str).apply(lambda row: " | ".join(row), axis=1))
                text_data += "\n"
            return text_data

        # CSV
        elif mime_type == "text/csv":
            df = pd.read_csv(file_path)
            return "\n".join(df.astype(str).apply(lambda row: " | ".join(row), axis=1))

        # PowerPoint
        elif mime_type in ["application/vnd.ms-powerpoint",
                           "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            prs = Presentation(file_path)
            text_data = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_data += shape.text + "\n"
            return text_data

        # Text / JSON / Markdown
        elif mime_type and (mime_type.startswith("text/") or
                            mime_type in ["application/json", "text/markdown"]):
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
    except Exception as e:
        print(f"Text extraction error from {file_path}: {e}")
    return ""


# -------------------------------
# Main embedding extraction
# -------------------------------
# Only replacing the video/audio parts in extract_embedding
def extract_embedding(file_path: str):
    mime_type, _ = mimetypes.guess_type(file_path)
    text = extract_text(file_path, mime_type)
    if text.strip():
        emb = clip_model.encode(text, convert_to_numpy=True)
        return emb.tolist(), text

    # Image
    if mime_type and mime_type.startswith("image/"):
        img = Image.open(file_path).convert("RGB")
        emb = clip_model.encode(img, convert_to_numpy=True)
        ocr_text = ocr_text_from_image_easyocr(img)
        if ocr_text:
            text_emb = clip_model.encode(ocr_text, convert_to_numpy=True)
            emb = (emb + text_emb) / 2
        return emb.tolist(), ocr_text if ocr_text else ""

    # Audio (split into segments)
    if mime_type and mime_type.startswith("audio/"):
        import librosa
        y, sr = librosa.load(file_path, sr=16000, mono=True)
        chunk_size = sr * 5  # 5 seconds per segment
        embeddings = []
        previews = []
        for i in range(0, len(y), chunk_size):
            segment = y[i:i + chunk_size]
            if len(segment) == 0:
                continue
            # Convert to spectrogram
            S = librosa.feature.melspectrogram(y=segment, sr=sr)
            S_dB = librosa.power_to_db(S, ref=np.max)
            fig, ax = plt.subplots(figsize=(3, 3), dpi=80)
            librosa.display.specshow(S_dB, sr=sr, x_axis=None, y_axis=None, ax=ax, cmap="magma")
            ax.axis("off")
            buf = io.BytesIO()
            plt.savefig(buf, format="png", bbox_inches="tight", pad_inches=0)
            plt.close(fig)
            buf.seek(0)
            img = Image.open(buf).convert("RGB")
            emb = clip_model.encode(img, convert_to_numpy=True)
            embeddings.append(emb.tolist())
            previews.append(f"[Audio segment {i // chunk_size}]")
        if embeddings:
            return embeddings, previews
        else:
            return [[0.0] * 512], [""]

    # Video (split into frame + audio segments)
    if mime_type and mime_type.startswith("video/"):
        try:
            import math
            cap = cv2.VideoCapture(file_path)
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            segment_duration = 5  # seconds
            segment_frames = int(fps * segment_duration)
            frame_embeddings = []
            previews = []

            segment_idx = 0
            frames_in_segment = []
            extracted_content = ""

            for i in range(frame_count):
                ret, frame = cap.read()
                if not ret:
                    break
                frames_in_segment.append(frame)

                if len(frames_in_segment) >= segment_frames or i == frame_count - 1:
                    # process segment frames
                    seg_embs = []
                    seg_text = ""
                    for f in frames_in_segment:
                        img = Image.fromarray(f).convert("RGB")
                        emb = clip_model.encode(img, convert_to_numpy=True)
                        ocr_text = ocr_text_from_image_easyocr(img)
                        if ocr_text:
                            text_emb = clip_model.encode(ocr_text, convert_to_numpy=True)
                            emb = (emb + text_emb) / 2
                            seg_text += ocr_text + "\n"
                        seg_embs.append(emb)
                    # average segment embeddings
                    if seg_embs:
                        seg_emb_avg = np.mean(seg_embs, axis=0)
                        frame_embeddings.append(seg_emb_avg.tolist())
                        previews.append(seg_text[:200] if seg_text else f"[Video segment {segment_idx}]")
                    frames_in_segment = []
                    segment_idx += 1

            cap.release()

            # Audio transcription per segment
            video_clip = VideoFileClip(file_path)
            with tempfile.NamedTemporaryFile(suffix=".wav") as tmp_audio:
                video_clip.audio.write_audiofile(tmp_audio.name, logger=None)
                result = whisper_model.transcribe(tmp_audio.name)
                audio_text = result.get("text", "")
                if audio_text.strip() and frame_embeddings:
                    audio_emb = clip_model.encode(audio_text, convert_to_numpy=True)
                    # combine audio with first video segment embedding (or all)
                    frame_embeddings[0] = (frame_embeddings[0] + audio_emb) / 2
                    previews[0] += "\n" + audio_text[:200]

            if frame_embeddings:
                return frame_embeddings, previews
        except Exception as e:
            print(f"Video processing error: {e}")
        return [[0.0] * 512], [""]

    return [[0.0] * 512], [""]


# -------------------------------
# Chunking helpers
# -------------------------------
def chunk_text(text: str, max_words: int = 500):
    words = text.split()
    chunks = []
    for i in range(0, len(words), max_words):
        chunks.append(" ".join(words[i:i + max_words]))
    return chunks


def extract_embeddings_with_chunks(file_path: str):
    full_embedding, full_text = extract_embedding(file_path)
    chunks = []

    if full_text.strip():
        text_chunks = chunk_text(full_text, max_words=500)
        for idx, chunk in enumerate(text_chunks):
            emb = clip_model.encode(chunk, convert_to_numpy=True).tolist()
            chunks.append({
                "chunk_embedding": emb,
                "chunk_content": chunk,
                "chunk_index": idx
            })

    return full_embedding, full_text, chunks
